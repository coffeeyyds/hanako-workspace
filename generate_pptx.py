#!/usr/bin/env python3
"""生成增强版 SpaceX IPO 制度性套利分析 PPT —— 含图表、表格、动效"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from lxml import etree
import os
import copy

# ── 配色 ──
BG_DARK  = RGBColor(0x0D, 0x0D, 0x12)
BG_SLIDE = RGBColor(0x14, 0x14, 0x1E)
ACCENT_GOLD  = RGBColor(0xD4, 0xA8, 0x40)
ACCENT_RED   = RGBColor(0xE0, 0x4B, 0x4B)
ACCENT_BLUE  = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_GREEN = RGBColor(0x3C, 0xB3, 0x71)
ACCENT_ORANGE= RGBColor(0xF0, 0x8C, 0x2E)
WHITE    = RGBColor(0xF0, 0xF0, 0xF0)
GRAY     = RGBColor(0x99, 0x99, 0xAA)
DARK_GRAY= RGBColor(0x55, 0x55, 0x66)
LIGHT_GRAY=RGBColor(0xBB, 0xBB, 0xCC)
CHART_BG = RGBColor(0x1A, 0x1A, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 全局字体 ──
FONT = "Microsoft YaHei"
FONT_EN = "Segoe UI"

# ═══════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════

def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name=FONT,
                  space_before=0, space_after=4, line_spacing=1.4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = Pt(line_spacing * font_size)
    return p

def add_section_title(slide, num, title, top=0.5):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(top), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_GOLD; circle.line.fill.background()
    ct = circle.text_frame; ct.paragraphs[0].text = str(num)
    ct.paragraphs[0].font.size = Pt(16); ct.paragraphs[0].font.color.rgb = BG_DARK
    ct.paragraphs[0].font.bold = True; ct.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_textbox(slide, 1.5, top-0.02, 8, 0.6, title, font_size=26, color=ACCENT_GOLD, bold=True)

def add_divider(slide, top):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Pt(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = DARK_GRAY; line.line.fill.background()

def add_bullet(tf, text, indent=0, font_size=16, color=WHITE, bold=False, bullet="•"):
    p = tf.add_paragraph()
    p.text = f"{'    '*indent}{bullet} {text}"
    p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = FONT; p.space_after = Pt(6); p.line_spacing = Pt(1.4*font_size)
    return p

def add_kv(tf, key, val, font_size=16, kc=GRAY, vc=WHITE, vb=False):
    p = tf.add_paragraph(); p.font.name = FONT; p.space_after = Pt(4); p.line_spacing = Pt(1.5*font_size)
    rk = p.add_run(); rk.text = f"{key}："; rk.font.size = Pt(font_size); rk.font.color.rgb = kc
    rv = p.add_run(); rv.text = val; rv.font.size = Pt(font_size); rv.font.color.rgb = vc; rv.font.bold = vb
    return p

def add_hl_box(slide, left, top, w, h, text, bg=ACCENT_RED, tc=WHITE, fs=20):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = bg; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs); p.font.color.rgb = tc
    p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    return tf

def add_table(slide, left, top, w, h, rows, cols, data, col_widths=None,
              header_bg=ACCENT_BLUE, body_bg=RGBColor(0x1E,0x1E,0x2E), font_size=14):
    """data: list of lists, first row = header"""
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(w), Inches(h))
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            p.font.name = FONT
            p.font.size = Pt(font_size)
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
            else:
                p.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = body_bg if r%2==0 else RGBColor(0x24,0x24,0x34)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl_shape

# ── 图表辅助 ──
def style_chart(chart, has_legend=True, legend_pos=XL_LEGEND_POSITION.BOTTOM):
    """统一美化图表：深色背景、白色文字、去掉边框"""
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = GRAY
    # 图表区
    chart_area = chart.chart_style
    # 绘图区
    plot = chart.plots[0]
    plot.gap_width = 80
    # 类别轴 & 数值轴
    try:
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(11)
        cat_axis.tick_labels.font.color.rgb = GRAY
        cat_axis.format.line.color.rgb = DARK_GRAY
    except: pass
    try:
        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = Pt(11)
        val_axis.tick_labels.font.color.rgb = GRAY
        val_axis.format.line.color.rgb = DARK_GRAY
        val_axis.major_gridlines.format.line.color.rgb = RGBColor(0x2A,0x2A,0x3A)
    except: pass

def add_data_labels(chart, font_size=10, color=WHITE):
    """给图表系列添加数据标签"""
    plot = chart.plots[0]
    for series in plot.series:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(font_size)
        series.data_labels.font.color.rgb = color
        series.data_labels.show_value = True

# ═══════════════════════════════════════════
# SLIDE 1: 封面
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 1.5, 2.0, 10.3, 1.5, "史无前例的制度性套利", font_size=48, color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.6, 10.3, 1.0, "SpaceX IPO × 被动投资的致命漏洞", font_size=28, color=ACCENT_GOLD)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.8), Inches(2.5), Pt(2.5))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_GOLD; line.line.fill.background()
add_textbox(slide, 1.5, 5.1, 5, 0.6, "2026 年 5 月", font_size=18, color=GRAY)

# ═══════════════════════════════════════════
# SLIDE 2: 目录
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_textbox(slide, 0.8, 0.6, 5, 0.8, "目录", font_size=36, color=WHITE, bold=True)
add_divider(slide, 1.5)
toc = [
    ("01","暴风雨前夜","三大巨头扎堆 IPO"),
    ("02","被动投资的庞大帝国","规模、逻辑与致命缺陷"),
    ("03","历史的镜子","特斯拉、日本小盘股、ARM"),
    ("04","SpaceX 的完美风暴","4% 流通股 × 三倍加权 × 15 天窗口"),
    ("05","赢家与输家","谁获益，谁埋单"),
    ("06","结语","非对称风险与监管真空"),
]
for i,(num,title,desc) in enumerate(toc):
    y=1.8+i*0.85
    add_textbox(slide,1.2,y,0.8,0.5,num,font_size=28,color=ACCENT_GOLD,bold=True)
    add_textbox(slide,2.2,y,4,0.5,title,font_size=22,color=WHITE,bold=True)
    add_textbox(slide,6.5,y+0.05,6,0.5,desc,font_size=16,color=GRAY)

# ═══════════════════════════════════════════
# SLIDE 3: 三大巨头 IPO 概览 + 柱状图
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"01","暴风雨前夜：三大巨头扎堆 IPO")
add_divider(slide,1.2)

# 左侧：三大公司卡片
cards = [
    ("SpaceX","估值 ~1.75 万亿美元","预计 6 月下旬 IPO\n史上最大规模",ACCENT_GOLD),
    ("OpenAI","超级独角兽","AGI 赛道领跑者\nChatGPT 母公司",ACCENT_BLUE),
    ("Anthropic","超级独角兽","安全 AI 旗舰\nClaude 背后团队",ACCENT_GREEN),
]
for i,(name,val,desc,clr) in enumerate(cards):
    x=0.8+i*2.9
    box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.6),Inches(2.6),Inches(2.4))
    box.fill.solid();box.fill.fore_color.rgb=RGBColor(0x1E,0x1E,0x2E);box.line.color.rgb=clr;box.line.width=Pt(1.5)
    tf=box.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=name;p.font.size=Pt(28);p.font.color.rgb=clr;p.font.bold=True;p.font.name=FONT;p.alignment=PP_ALIGN.CENTER
    add_paragraph(tf,val,font_size=16,color=WHITE,alignment=PP_ALIGN.CENTER)
    add_paragraph(tf,desc,font_size=13,color=GRAY,alignment=PP_ALIGN.CENTER)

# 右侧：柱状图 - IPO 规模对比
chart_data = CategoryChartData()
chart_data.categories = ['SpaceX\n(预计)', 'OpenAI\n(预计)', 'Anthropic\n(预计)', 'ARM\n(2023)', 'Rivian\n(2021)', 'Uber\n(2019)']
chart_data.add_series('估值/募资（亿美元）', [17500, 1500, 600, 545, 665, 824])
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(9.5), Inches(1.6), Inches(3.3), Inches(3.0), chart_data
).chart
style_chart(chart, has_legend=False)
chart.plots[0].gap_width = 60
add_data_labels(chart, font_size=8)
# 给 SpaceX 的柱子特殊颜色
series = chart.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = ACCENT_GOLD
# 数轴
try:
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.color.rgb = GRAY
except: pass
try:
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = GRAY
except: pass

add_textbox(slide, 9.5, 4.8, 3.3, 0.4, "▲ SpaceX 估值远超其他 IPO", font_size=11, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_hl_box(slide,0.8,4.4,8.2,2.5,
    "每一个单独拿出来都能成为年度 IPO 事件\n三家愣要挤在一起上市\n\n当大家还在惊叹万亿美元市值时\n这些未上市公司的估值已经破万亿了",
    bg=RGBColor(0x2A,0x2A,0x35),tc=ACCENT_GOLD,fs=18)

# ═══════════════════════════════════════════
# SLIDE 4: 指数公司的两难
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"01","指数公司的两难")
add_divider(slide,1.2)

tf=add_textbox(slide,0.8,1.5,5.8,1.0,"如果 SpaceX、OpenAI、Anthropic 都不在指数里…",font_size=24,color=WHITE,bold=True)
tf2=add_textbox(slide,0.8,2.5,5.8,0.8,"受损的不是这几家的信誉，而是「指数是否有效」的争议",font_size=20,color=ACCENT_RED,bold=True)
tf3=add_textbox(slide,0.8,3.4,5.8,3.5,"",font_size=16,color=WHITE)
add_bullet(tf3,"纳斯达克 100 → 须反映科技股综合表现")
add_bullet(tf3,"标普 500 → 须反映美股主要公司估值")
add_bullet(tf3,"SpaceX = 商业航天一骑绝尘，代表人类先进科技")
add_bullet(tf3,"OpenAI / Anthropic = AGI 赛道领跑者")
add_bullet(tf3,"不纳入 → 指数失去代表性 → 指数公司的存在意义受质疑")

# 右侧流程图：指数公司困境
add_hl_box(slide,7.5,1.5,5.0,5.0,
    "指数公司的逻辑死结\n\n不纳入\n↓\n指数失效\n失去公信力\n\n纳入\n↓\n必须改规则\n↓\n埋下套利隐患",
    bg=RGBColor(0x2A,0x2A,0x35),tc=ACCENT_BLUE,fs=18)

# ═══════════════════════════════════════════
# SLIDE 5: 规则对比表格
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"01","规则为谁而改？ —— 指数纳入规则对比")
add_divider(slide,1.2)

add_table(slide,0.8,1.5,11.7,2.8,4,5,
    [["","标普 500（现行）","标普 500（新提案）","纳斯达克 100（现行）","纳斯达克 100（新提案）"],
     ["纳入等待期","上市满 12 个月","6 个月或更短","3 个月","仅 15 天"],
     ["门槛条件","盈利要求 + 市值","市值 > $2000 亿\n即可走绿色通道","标准市值门槛","取消 10% 流通量\n硬性限制"],
     ["额外加权","无","无","无","流通股 < 20% →\n权重 × 3 倍"]],
    col_widths=[2.0,2.4,2.4,2.4,2.5],
    header_bg=ACCENT_BLUE,font_size=14)

# 底部
add_hl_box(slide,0.8,4.8,11.7,2.2,
    "这是摆明了为 SpaceX / OpenAI / Anthropic 量身定制的规则修改\n标普先出招，纳斯达克动作更激进，直接把等待期压到 15 天",
    bg=RGBColor(0x2A,0x2A,0x35),tc=ACCENT_GOLD,fs=18)

# ═══════════════════════════════════════════
# SLIDE 6: 被动投资的庞大帝国 + 饼图
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"02","被动投资的庞大帝国")
add_divider(slide,1.2)

tf=add_textbox(slide,0.8,1.5,5.8,1.0,"为什么被动投资大行其道？",font_size=24,color=WHITE,bold=True)
tf2=add_textbox(slide,0.8,2.5,5.8,4.0,"",font_size=16,color=WHITE)
add_bullet(tf2,"散户高抛低吸 = 稳定亏手续费")
add_bullet(tf2,'主动基金经理 = "穿着西装的猴子扔飞镖"')
add_bullet(tf2,"长期跑赢大盘的经理寥寥无几 → 幸存者偏差")
add_bullet(tf2,"被动投资：管理费极低、无利润分成、回报丰厚")
add_bullet(tf2,"已成主流投资方式 → 美股 ETF AUM 超 10 万亿美元")

# 饼图：被动 vs 主动
pie_data = CategoryChartData()
pie_data.categories = ['指数被动基金', '主动管理基金', '对冲基金等']
pie_data.add_series('AUM占比', [42, 38, 20])
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(8.0), Inches(1.5), Inches(4.5), Inches(3.5), pie_data
)
chart = chart_shape.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = GRAY
plot = chart.plots[0]
# 配色
colors_pie = [ACCENT_BLUE, ACCENT_GOLD, DARK_GRAY]
for i, color in enumerate(colors_pie):
    point = plot.series[0].points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color
# 数据标签
plot.series[0].has_data_labels = True
plot.series[0].data_labels.font.size = Pt(11)
plot.series[0].data_labels.font.color.rgb = WHITE
plot.series[0].data_labels.show_percentage = True
plot.series[0].data_labels.show_category_name = True

add_textbox(slide,8.0,5.2,4.5,0.4,"▲ 美股市场资金结构（示意）",font_size=11,color=GRAY,alignment=PP_ALIGN.CENTER)

add_hl_box(slide,0.8,5.8,6.8,1.2,
    "VOO 单只 ETF 规模突破万亿美元，堪称奇观",
    bg=ACCENT_GOLD,tc=BG_DARK,fs=16)

# ═══════════════════════════════════════════
# SLIDE 7: 跟踪误差
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"02","跟踪误差：管钥匙的丫鬟")
add_divider(slide,1.2)

tf=add_textbox(slide,0.8,1.5,5.8,5.5,"",font_size=16)
add_paragraph(tf,"指数公司 v.s. ETF 基金",font_size=24,color=WHITE,bold=True)
add_bullet(tf,"指数公司只编指数 → 纳入/剔除瞬间完成 = 「真空球形鸡」")
add_bullet(tf,"ETF 基金必须在真实市场买卖 → 手续费、冲击成本、流动性限制")
add_bullet(tf,"核心 KPI = 跟踪误差（越小越好），代价由投资人管理费埋单")
add_paragraph(tf,"",font_size=6)
add_paragraph(tf,"核心矛盾",font_size=22,color=ACCENT_RED,bold=True)
add_bullet(tf,"基金经理只是「管钥匙的丫鬟」",color=WHITE,font_size=16)
add_bullet(tf,"实际主人 = 普通打工人的退休金 + 养老基金",color=WHITE,font_size=16)
add_bullet(tf,"丫鬟必须按主人的规则不计代价执行",color=WHITE,font_size=16)

add_hl_box(slide,7.5,1.8,5.0,5.0,
    "这只「活鸡」\n已经巨大到\n对资本市场有\n举足轻重的影响\n\n且已经出过不少事了",
    bg=ACCENT_RED,tc=WHITE,fs=22)

# ═══════════════════════════════════════════
# SLIDE 8: 特斯拉案例 + 时间线 + 价格走势示意
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"03","历史教训① 特斯拉纳入 S&P 500（2020）")
add_divider(slide,1.2)

# 时间线条形图：TSLA 纳入前后价格变化示意
chart_data = CategoryChartData()
chart_data.categories = ['2020.07\n布局期', '2020.09\n预期升温', '2020.11\n宣布纳入', '2020.12.21\n纳入日', '2021.01\n高位', '2021.03\n回调']
chart_data.add_series('TSLA 股价走势（示意）', [280, 380, 540, 695, 880, 560])
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5), chart_data
)
chart = chart_shape.chart
style_chart(chart, has_legend=False)
chart.plots[0].gap_width = 50
add_data_labels(chart, font_size=9)
series = chart.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = ACCENT_BLUE
# 纳入日柱子标红
pt = series.points[3]
pt.format.fill.solid()
pt.format.fill.fore_color.rgb = ACCENT_RED

add_textbox(slide,7.0,5.2,5.5,0.4,"▲ TSLA 纳入 S&P 500 前后股价示意",font_size=11,color=GRAY,alignment=PP_ALIGN.CENTER)

# 左侧文字
tf=add_textbox(slide,0.8,1.5,5.8,5.0,"",font_size=16)
add_paragraph(tf,"时间线",font_size=22,color=WHITE,bold=True)
add_bullet(tf,"2020 夏：主动基金大量囤积 TSLA → 提前布局")
add_bullet(tf,"11 月：宣布纳入 → 几周内飙升 50%+")
add_bullet(tf,'12.21「The Big Bang」：被动基金必须买入 ≈ 流通盘 15%~20%',color=ACCENT_RED,bold=True)
add_bullet(tf,"当日成交突破 1500 亿美元 → 以历史最高溢价成交")
add_bullet(tf,"2021：剧烈回调 → 被动投资者平均买贵 30%+",color=ACCENT_RED)

add_hl_box(slide,0.8,5.8,11.7,1.2,
    "主动基金提前囤货抬价 → 被动基金不计代价接盘 → 合法打劫普通投资人 → 之后无人追责",
    bg=RGBColor(0x3A,0x1A,0x1A),tc=ACCENT_RED,fs=16)

# ═══════════════════════════════════════════
# SLIDE 9: 日本小盘股
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"03","历史教训② 日本小盘股流动性真空（2024）")
add_divider(slide,1.2)

tf=add_textbox(slide,0.8,1.5,5.8,5.0,"",font_size=16)
add_paragraph(tf,"极端案例",font_size=24,color=WHITE,bold=True)
add_bullet(tf,"指数规则调整 → 被动基金被迫买入冷门小盘股")
add_bullet(tf,"买入量 > 该股连续 30 天日均成交量 × 1000 倍",color=ACCENT_RED,bold=True,font_size=18)
add_bullet(tf,"连续数天无量涨停")
add_bullet(tf,"指数公司被迫撤回指令")
add_paragraph(tf,"",font_size=6)
add_paragraph(tf,"致命推论",font_size=24,color=WHITE,bold=True)
add_bullet(tf,"如果被动 ETF 买入量 > 流通总量",color=WHITE)
add_bullet(tf,"理论上股价 → 无穷大",color=ACCENT_RED,bold=True,font_size=18)
add_bullet(tf,"资金有限，但漏洞确实存在，且很大",color=WHITE)

# 右侧：流动性对比柱状图
chart_data2 = CategoryChartData()
chart_data2.categories = ['正常股票', '冷门小盘股\n（日本案例）', 'SpaceX\n（4%流通股）']
chart_data2.add_series('被动基金买入 / 日均成交量（倍）', [2, 1000, 50])
chart_shape2 = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5), chart_data2
)
chart2 = chart_shape2.chart
style_chart(chart2, has_legend=False)
chart2.plots[0].gap_width = 60
add_data_labels(chart2, font_size=10)
s = chart2.series[0]
s.format.fill.solid()
s.format.fill.fore_color.rgb = ACCENT_RED

add_textbox(slide,7.5,6.3,5.0,0.5,"▲ 买入需求 vs 流动性的失衡程度",font_size=11,color=GRAY,alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 10: ARM 连环套利 + 流程图
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"03","历史教训③ ARM 的连环套利")
add_divider(slide,1.2)

# 三步卡片
steps = [
    ("① 惜售流通股","仅释放 ~10%\n软银现金流紧张\nARM 是手上王牌",ACCENT_BLUE),
    ("② Gamma Squeeze","少量 CALL 期权\n→ 做市商被迫买入对冲\n→ 流通盘被抢 → 跳空暴涨",ACCENT_ORANGE),
    ("③ 质押套现","90% 股票拿去银行质押\n少量流通股拉高总市值\n→ 天量信贷进账",ACCENT_GREEN),
]
for i,(title,desc,clr) in enumerate(steps):
    x=0.8+i*4.2
    # 箭头（非最后一个）
    if i<2:
        arrow=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x+3.2),Inches(3.5),Inches(0.8),Inches(0.4))
        arrow.fill.solid();arrow.fill.fore_color.rgb=DARK_GRAY;arrow.line.fill.background()
    box=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(1.8),Inches(3.2),Inches(3.5))
    box.fill.solid();box.fill.fore_color.rgb=RGBColor(0x1E,0x1E,0x2E);box.line.color.rgb=clr;box.line.width=Pt(1.5)
    tf=box.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=title;p.font.size=Pt(22);p.font.color.rgb=clr;p.font.bold=True;p.font.name=FONT;p.alignment=PP_ALIGN.CENTER
    add_paragraph(tf,desc,font_size=14,color=WHITE,alignment=PP_ALIGN.CENTER)

add_hl_box(slide,0.8,5.8,11.7,1.2,
    "学好不容易，学坏一出溜 → 这个套路能在场外造成次级效应，被所有人看到了",
    bg=RGBColor(0x3A,0x1A,0x1A),tc=ACCENT_GOLD,fs=16)

# ═══════════════════════════════════════════
# SLIDE 11: SpaceX vs ARM 流通股对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"04","完美风暴：仅 4% 流通股")
add_divider(slide,1.2)

# 流通股比例对比柱状图
chart_data = CategoryChartData()
chart_data.categories = ['SpaceX\n(计划)', 'ARM\n(2023)', '正常 IPO\n典型范围']
chart_data.add_series('流通股比例 (%)', [4.3, 10, 30])
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5), chart_data
)
chart = chart_shape.chart
style_chart(chart, has_legend=False)
add_data_labels(chart, font_size=11)
s = chart.series[0]
s.format.fill.solid()
s.format.fill.fore_color.rgb = ACCENT_BLUE
s.points[0].format.fill.solid()
s.points[0].format.fill.fore_color.rgb = ACCENT_RED

# 左侧对比
tf=add_textbox(slide,0.8,1.5,5.8,2.0,"",font_size=16)
add_paragraph(tf,"ARM（前车之鉴）",font_size=26,color=ACCENT_BLUE,bold=True)
add_kv(tf,"流通股比例","~10%",font_size=18,vc=WHITE)
add_kv(tf,"后果","股价翻倍 + 天量质押套现",font_size=16,vc=ACCENT_RED)

tf2=add_textbox(slide,0.8,3.6,5.8,2.0,"",font_size=16)
add_paragraph(tf2,"SpaceX（升级版）",font_size=26,color=ACCENT_GOLD,bold=True)
add_kv(tf2,"流通股比例","仅 ~4.3%",font_size=22,vc=ACCENT_RED,vb=True)
add_kv(tf2,"IPO 释放金额","~$750 亿",font_size=18,vc=WHITE)
add_kv(tf2,"总市值","~$1.75 万亿",font_size=18,vc=WHITE)

add_hl_box(slide,0.8,5.8,11.7,1.2,
    "比 ARM 更极端：仅 4% 流通 + 15 天窗口 + 三倍加权 = 人为制造的流动性完美风暴",
    bg=ACCENT_RED,tc=WHITE,fs=18)

# ═══════════════════════════════════════════
# SLIDE 12: 三倍加权 - 权重对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"04","三倍加权：纳斯达克的神操作")
add_divider(slide,1.2)

# 权重变化柱状图
chart_data = CategoryChartData()
chart_data.categories = ['旧规则\n（流通市值）', '新规则\n（×3 倍）']
chart_data.add_series('SpaceX 在 NDX 中的权重市值 ($亿)', [750, 2250])
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5), chart_data
)
chart = chart_shape.chart
style_chart(chart, has_legend=False)
add_data_labels(chart, font_size=12)
s = chart.series[0]
s.format.fill.solid()
s.format.fill.fore_color.rgb = ACCENT_BLUE
s.points[1].format.fill.solid()
s.points[1].format.fill.fore_color.rgb = ACCENT_RED

add_textbox(slide,7.0,5.2,5.5,0.4,"▲ SpaceX 指数权重市值：$750 亿 → $2250 亿",font_size=11,color=ACCENT_GOLD,alignment=PP_ALIGN.CENTER)

# 左侧规则
tf=add_textbox(slide,0.8,1.5,5.8,3.0,"",font_size=16)
add_paragraph(tf,"纳斯达克新规则（针对前 40 大 IPO）",font_size=22,color=ACCENT_GOLD,bold=True)
add_bullet(tf,"取消「流通量 < 10% 不准进」的硬性门槛")
add_bullet(tf,"流通股 < 20% → 权重按 3 倍计算",color=ACCENT_RED,bold=True)
add_bullet(tf,"SpaceX 排名跃升至 NDX 前十")
add_bullet(tf,"权重可能超过 Netflix 或 Palantir")

add_hl_box(slide,0.8,4.6,5.8,2.5,
    "指数公司的尴尬：\n按真实流通量 → 存在感太低\n按 100% → 股票不够买\n\n「三倍」是折中方案\n却恰好制造了最大套利空间",
    bg=RGBColor(0x2A,0x2A,0x35),tc=ACCENT_GOLD,fs=16)

# ═══════════════════════════════════════════
# SLIDE 13: 15天窗口期的数学暴力
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"04","15 天窗口期的数学暴力")
add_divider(slide,1.2)

# 供需缺口对比柱状图
chart_data = CategoryChartData()
chart_data.categories = ['SpaceX\n流通盘', 'NDX\n需配置', 'SPX\n需配置']
chart_data.add_series('金额（亿美元）', [750, 2250, 5000])
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5), chart_data
)
chart = chart_shape.chart
style_chart(chart, has_legend=False)
add_data_labels(chart, font_size=10)
s = chart.series[0]
s.format.fill.solid()
s.format.fill.fore_color.rgb = ACCENT_BLUE
s.points[0].format.fill.solid()
s.points[0].format.fill.fore_color.rgb = ACCENT_GREEN
s.points[1].format.fill.solid()
s.points[1].format.fill.fore_color.rgb = ACCENT_RED
s.points[2].format.fill.solid()
s.points[2].format.fill.fore_color.rgb = ACCENT_ORANGE

add_textbox(slide,7.0,5.2,5.5,0.8,"▲ 绿=供给 红/橙=需求\n流通盘总量 ≈ 买盘需求量级",font_size=11,color=GRAY,alignment=PP_ALIGN.CENTER)

tf=add_textbox(slide,0.8,1.5,5.8,5.0,"",font_size=16)
add_paragraph(tf,"强制买入 = 不计代价",font_size=22,color=WHITE,bold=True)
add_bullet(tf,"NDX 挂钩产品必须配置 $2250 亿 → 占指数 0.7%~0.8%")
add_bullet(tf,"SPX 类资产池（>20 万亿美元）也必须买几百亿")
add_bullet(tf,"总流通盘 ≈ 总买盘需求量级",color=ACCENT_RED,bold=True)
add_bullet(tf,"15 天内套利者必须抢先布局 → 抢 4% 的流通股",color=ACCENT_RED)
add_paragraph(tf,"",font_size=6)
add_paragraph(tf,"市值只是一个估算值，并不真的代表公司价值\n但被动基金必须为这个估值买单",font_size=18,color=ACCENT_GOLD)

# ═══════════════════════════════════════════
# SLIDE 14: 马斯克的后手
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"04","马斯克的两手后招")
add_divider(slide,1.2)

tf=add_textbox(slide,0.8,1.5,5.5,2.5,"",font_size=16)
add_paragraph(tf,"招数一：二级增发",font_size=26,color=ACCENT_GOLD,bold=True)
add_bullet(tf,"股价已被推上天 → 再发股票融资合理合法")
add_bullet(tf,'"既然抢着要，我给你们，我是大善人"')
add_bullet(tf,"同时给市场一个台阶下，避免被国会传唤")

tf2=add_textbox(slide,7.5,1.5,5.0,2.5,"",font_size=16)
add_paragraph(tf2,"招数二：质押贷款",font_size=26,color=ACCENT_GOLD,bold=True)
add_bullet(tf2,"公司总市值已被拉到天价")
add_bullet(tf2,"总市值 × 抵押率 → 天量融资",color=ACCENT_RED,bold=True)
add_bullet(tf2,"银行 or 私募信贷公司承接")

# 底部流程图
add_hl_box(slide,0.8,4.5,11.7,2.5,
    "完整路径\n\n极小流通股(4%) → 三倍加权虚增权重 → 15天窗口逼迫抢筹\n→ 被动基金不计代价买入 → 股价火箭 → 增发+质押双线收割\n→ 银行/私募承接风险 → 出事了就「没料到」→ 非对称风险套利闭环",
    bg=RGBColor(0x3A,0x1A,0x1A),tc=ACCENT_GOLD,fs=16)

# ═══════════════════════════════════════════
# SLIDE 15: 赢家 vs 输家 对照表
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_section_title(slide,"05","谁获益，谁埋单？")
add_divider(slide,1.2)

# 对照表
add_table(slide,0.8,1.5,5.5,4.5,5,2,
    [["⚠ 埋单的人",""],
     ["指数被动投资散户","相信「一劳永逸」的普通人"],
     ["退休金 / 养老基金","指望退休金养老的群体"],
     ["被动 ETF 持有人","被合法打劫的羊群"],
     ["","又一次，管钥匙的丫鬟\n用主人的钱把主人宰了一顿"]],
    col_widths=[2.5,3.0],
    header_bg=ACCENT_RED,font_size=15)

add_table(slide,7.0,1.5,5.5,4.5,5,2,
    [["✓ 获益的人",""],
     ["马斯克本人","万亿级身家暴涨"],
     ["SpaceX 早期投资人","私募轮回报爆炸"],
     ["15 天窗口套利者","利用 Gamma Squeeze"],
     ["","万亿级别的资本盛筵"]],
    col_widths=[2.5,3.0],
    header_bg=ACCENT_GREEN,font_size=15)

add_textbox(slide,0.8,6.5,11.5,0.8,
    "监管部门？以后还要走玻璃旋转门去投行上班的，非得阻碍未来老板和客户发财吗？",
    font_size=15,color=GRAY,alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 16: 结语
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide,1.5,1.5,10.3,1.2,"非对称风险套利",font_size=48,color=WHITE,bold=True)
add_textbox(slide,1.5,3.0,10.3,3.0,
    "赢家拿走全部收益，输家承担所有风险。\n\n规则由最有能力利用规则的人来书写。\n\n也许我的推演完全都是错的，\n未来是严格不可预测的。\n\n且来看这场好戏。",
    font_size=24,color=GRAY)

line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.5),Inches(6.2),Inches(2.5),Pt(2.5))
line.fill.solid();line.fill.fore_color.rgb=ACCENT_GOLD;line.line.fill.background()

add_textbox(slide,1.5,6.5,5,0.5,"谢谢阅读",font_size=20,color=GRAY)

# ═══════════════════════════════════════════
# 添加幻灯片切换动画（淡入淡出）
# ═══════════════════════════════════════════
nsmap = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
}

def add_fade_transition(slide, duration="slow"):
    """给幻灯片添加淡入淡出过渡动画"""
    sld = slide._element
    # 确保 transition 节点存在
    trans = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    if trans is None:
        # 在 sld 的最后一个子元素前插入
        cSld = sld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
        trans = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        sld.insert(list(sld).index(cSld) + 1, trans)
    
    # 添加淡入淡出效果
    fade = etree.SubElement(trans, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
    # 持续时间
    if duration == "slow":
        fade.set('dur', '800')  # ms
    elif duration == "medium":
        fade.set('dur', '500')
    else:
        fade.set('dur', '300')

# 给所有幻灯片添加淡入淡出切换
for slide in prs.slides:
    add_fade_transition(slide, "medium")

# ── 保存 ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SpaceX_IPO_制度性套利分析.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print("Includes: charts, tables, and fade transitions on all 16 slides.")
